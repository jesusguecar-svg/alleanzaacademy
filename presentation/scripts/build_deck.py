#!/usr/bin/env python3
"""Deck ejecutivo Alleanza Academy (ES) v2 — con logo + curso, competencia y pricing."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PURPLE  = RGBColor(0x7C,0x46,0xA6)
PURPLE2 = RGBColor(0x99,0x58,0xD7)
LIGHT   = RGBColor(0xEA,0xEC,0xF3)
NAVY    = RGBColor(0x06,0x13,0x30)
NAVY2   = RGBColor(0x12,0x20,0x42)
WHITE   = RGBColor(0xFF,0xFF,0xFF)
GREY    = RGBColor(0x5B,0x60,0x70)
GREEN   = RGBColor(0x2E,0x7D,0x5B)
FONT    = "Inter"
B = "/tmp/claude-0/-home-user-alleanzaacademy/147edd4c-26bd-5fd6-a68b-8407a942b324/scratchpad/"
LOGO_C, LOGO_W = B+"alleanza_logo.png", B+"alleanza_logo_white.png"

prs = Presentation()
prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
BLANK=prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(BLANK)
def bg(s,c):
    s.background.fill.solid(); s.background.fill.fore_color.rgb=c
def rect(s,x,y,w,h,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background(); sh.shadow.inherit=False
    return sh
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sa=6,ls=1.05):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(sa); p.line_spacing=ls
        for (t,sz,c,b,*r) in (para if isinstance(para,list) else [para]):
            run=p.add_run(); run.text=t; run.font.name=FONT; run.font.size=Pt(sz)
            run.font.color.rgb=c; run.font.bold=b
            if r: run.font.italic=r[0]
    return tb
def logo(s,dark=False):
    s.shapes.add_picture(LOGO_W if dark else LOGO_C, Inches(11.15), Inches(0.42), height=Inches(0.52))
def accent(s): rect(s,0,0,Inches(0.18),SH,PURPLE)
def kicker(s,t,c=PURPLE): txt(s,Inches(0.9),Inches(0.62),Inches(8),Inches(0.4),[[(t.upper(),13,c,True)]])
def title(s,t,c=NAVY,y=1.0,size=31): txt(s,Inches(0.9),Inches(y),Inches(10),Inches(1.0),[[(t,size,c,True)]],ls=1.0)
def footer(s,n,dark=False):
    c=LIGHT if dark else GREY
    txt(s,Inches(0.9),Inches(7.04),Inches(9),Inches(0.35),[[("Alleanza Academy · Alineación estratégica previa al lanzamiento",9,c,False)]])
    txt(s,Inches(12.2),Inches(7.04),Inches(0.7),Inches(0.35),[[(str(n),9,c,False)]],align=PP_ALIGN.RIGHT)
def card(s,x,y,w,h,head,body,hc=PURPLE,dark=False):
    rect(s,x,y,w,h,NAVY2 if dark else WHITE); rect(s,x,y,w,Inches(0.09),hc)
    txt(s,x+Inches(0.22),y+Inches(0.26),w-Inches(0.44),Inches(0.7),[[(head,14.5,WHITE if dark else NAVY,True)]])
    txt(s,x+Inches(0.22),y+Inches(0.92),w-Inches(0.44),h-Inches(1.05),[[(body,11.5,LIGHT if dark else GREY,False)]],ls=1.12)
def bullets(s,x,y,w,h,items,size=14,gap=10,dark=False):
    runs=[[("●  ",size,PURPLE2 if dark else PURPLE,True),(it,size,WHITE if dark else NAVY,False)] for it in items]
    txt(s,x,y,w,h,runs,sa=gap,ls=1.08)

# ============ 1. PORTADA ============
s=slide(); bg(s,NAVY); rect(s,0,0,Inches(0.22),SH,PURPLE)
s.shapes.add_picture(LOGO_W, Inches(0.85), Inches(1.55), width=Inches(6.3))
rect(s,Inches(0.95),Inches(4.25),Inches(3.2),Inches(0.05),PURPLE)
txt(s,Inches(0.95),Inches(4.5),Inches(11.5),Inches(0.6),
    [[("Preparación para el examen de licencia de seguros de Salud y Vida",17,LIGHT,False)]])
txt(s,Inches(0.95),Inches(5.35),Inches(11.5),Inches(1.2),
    [[("Reunión de alineación estratégica previa al lanzamiento",15,WHITE,True)],
     [("Participantes: Jesús Esteban Guerra · Guillermo · Cynthia",13,LIGHT,False)],
     [("23 de junio de 2026",13,GREY,False)]],sa=5)

# ============ 2. POR QUÉ / VISIÓN ============
s=slide(); bg(s,WHITE); accent(s); logo(s); kicker(s,"Por qué — la visión")
title(s,"Una academia al servicio de la visión integral de la organización")
txt(s,Inches(0.9),Inches(1.9),Inches(11.4),Inches(0.9),
    [[("El objetivo no es vender cursos: es ",15,NAVY,False),("resolver in-house cada necesidad del agente",15,PURPLE,True),
      (". El gran ingreso viene de la producción; la academia acelera que más agentes se licencien y produzcan, antes y mejor.",15,NAVY,False)]],ls=1.15)
y,w,h,g=3.05,3.75,3.05,0.18
card(s,Inches(0.9),Inches(y),Inches(w),Inches(h),"Todo dentro de casa",
     "El agente obtiene su licencia, certificación y formación continua sin salir de Alleanza. Menos fricción, más retención.")
card(s,Inches(0.9+w+g),Inches(y),Inches(w),Inches(h),"Acelera la producción",
     "Cada agente licenciado antes es un agente que produce antes. La academia alimenta el motor de ingresos real.")
card(s,Inches(0.9+2*(w+g)),Inches(y),Inches(w),Inches(h),"Autosostenible",
     "Se financia con las inscripciones: no es un centro de costo, es un recurso que se paga solo y crece con la organización.")
footer(s,2)

# ============ 3. EL CURSO ============
s=slide(); bg(s,WHITE); accent(s); logo(s); kicker(s,"El curso")
title(s,"Preparación real para aprobar el examen de Salud y Vida")
txt(s,Inches(0.9),Inches(1.85),Inches(6.1),Inches(1.5),
    [[("Nuestro propósito: ",15,NAVY,True),("ayudar al agente a pasar en su siguiente intento",15,PURPLE,True),
      (", sea el primero o no. No preparamos para lucir bien; preparamos para aprobar — y en español, con los términos clave en inglés.",15,NAVY,False)]],ls=1.15)
# núcleo: simulador
rect(s,Inches(0.9),Inches(3.35),Inches(6.0),Inches(2.85),PURPLE)
txt(s,Inches(1.2),Inches(3.6),Inches(5.5),Inches(0.5),[[("EL NÚCLEO DEL MÉTODO",13,WHITE,True)]])
txt(s,Inches(1.2),Inches(4.1),Inches(5.5),Inches(0.8),[[("Simulador de examen real",26,WHITE,True)]])
txt(s,Inches(1.2),Inches(4.95),Inches(5.4),Inches(1.2),
    [[("Recrea las condiciones reales del examen para que el agente llegue verdaderamente preparado. Todo lo demás es complemento de esto.",13,LIGHT,False)]],ls=1.18)
# accesorios
txt(s,Inches(7.4),Inches(1.95),Inches(5),Inches(0.5),[[("RECURSOS COMPLEMENTARIOS",13,PURPLE,True)]])
bullets(s,Inches(7.4),Inches(2.5),Inches(5.2),Inches(4.2),[
    "Clases en vivo por Zoom",
    "Lecciones en video interactivo",
    "Flashcards de repaso",
    "Audio / podcasts para estudiar en cualquier momento",
    "Material de lectura estructurado",
],size=14.5,gap=15)
footer(s,3)

# ============ 4. VENTAJA COMPETITIVA (tabla) ============
s=slide(); bg(s,WHITE); accent(s); logo(s); kicker(s,"Posicionamiento")
title(s,"Por qué Alleanza Academy gana donde importa: aprobar")
rows=[
 ["", "Alleanza Academy", "Kaplan / Xcel\n(nacionales)", "Instructores sueltos\n(María Papi · Aaron)"],
 ["Propósito","Que apruebes el examen","Cumplir el requisito de curso","Dar clase / repaso"],
 ["Idioma","Español, con términos clave en inglés","Solo en inglés","Español"],
 ["Simulador realista","Sí — núcleo del método","Genérico","Pocas diapositivas / no"],
 ["Formatos de estudio","Zoom + video + flashcards + podcast + lectura","Online autoguiado","Un solo formato"],
 ["Pre-licencia (req. FL)","Opcional — te acompañamos","Incluido","A veces, con fee aparte"],
 ["Inversión","$107 lanzamiento (valor $350)","Mayor","$100 (María) / mensual (Aaron)"],
]
from pptx.util import Inches as I
left,top,tw,th=I(0.9),I(1.82),I(11.5),I(4.55)
tbl=s.shapes.add_table(len(rows),4,left,top,tw,th).table
tbl.columns[0].width=I(2.3); tbl.columns[1].width=I(3.4); tbl.columns[2].width=I(2.8); tbl.columns[3].width=I(3.0)
for ri,row in enumerate(rows):
    tbl.rows[ri].height=I(0.5 if ri==0 else 0.62)
    for ci,val in enumerate(row):
        cell=tbl.cell(ri,ci); cell.margin_left=I(0.1); cell.margin_right=I(0.08)
        cell.margin_top=I(0.04); cell.margin_bottom=I(0.04); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        tf=cell.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.LEFT
        hl = (ri==2 and ci==1)   # fila Idioma, columna Alleanza -> destacada
        # fill
        if ri==0:
            cell.fill.solid(); cell.fill.fore_color.rgb=NAVY
        elif hl:
            cell.fill.solid(); cell.fill.fore_color.rgb=PURPLE
        elif ci==1:
            cell.fill.solid(); cell.fill.fore_color.rgb=LIGHT
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb=WHITE
        for j,line in enumerate(val.split("\n")):
            pp=p if j==0 else tf.add_paragraph(); pp.alignment=PP_ALIGN.LEFT
            r=pp.add_run(); r.text=line; r.font.name=FONT
            if ri==0:
                r.font.size=Pt(12.5); r.font.bold=True
                r.font.color.rgb=WHITE if ci!=1 else PURPLE2
            elif hl:
                r.font.size=Pt(11); r.font.bold=True; r.font.color.rgb=WHITE
            elif ci==0:
                r.font.size=Pt(11.5); r.font.bold=True; r.font.color.rgb=NAVY
            elif ci==1:
                r.font.size=Pt(11); r.font.bold=True; r.font.color.rgb=PURPLE
            else:
                r.font.size=Pt(10.5); r.font.bold=False; r.font.color.rgb=GREY
txt(s,Inches(0.9),Inches(6.55),Inches(11.4),Inches(0.4),
    [[("Honestidad: ",10.5,PURPLE,True),("no cumplimos el requisito de pre-licencia como Kaplan/Xcel — por eso lo ofrecemos como apoyo opcional. Nuestro foco es que el agente apruebe.",10.5,GREY,False,True)]])
footer(s,4)

# ============ 5. URGENCIA MIAMI + RUTAS ============
s=slide(); bg(s,WHITE); accent(s); logo(s); kicker(s,"La urgencia — equipos como Miami")
title(s,"Rutas según el tiempo de cada agente")
txt(s,Inches(0.9),Inches(1.85),Inches(11.3),Inches(0.8),
    [[("Equipos como Miami necesitan licenciar y capacitar agentes ya. Respondemos con rutas flexibles, no un calendario rígido — con el mismo estándar de calidad.",14.5,NAVY,False)]],ls=1.15)
rect(s,Inches(0.9),Inches(3.0),Inches(5.5),Inches(2.95),PURPLE)
txt(s,Inches(1.2),Inches(3.25),Inches(5),Inches(0.5),[[("RUTA INTENSIVA",15,WHITE,True)]])
txt(s,Inches(1.2),Inches(3.72),Inches(5),Inches(0.7),[[("≈ 2 semanas al examen",23,WHITE,True)]])
txt(s,Inches(1.2),Inches(4.5),Inches(4.9),Inches(1.3),[[("Para quienes deben presentar el examen con urgencia. Ritmo acelerado, enfocado 100% en aprobar.",13,LIGHT,False)]],ls=1.18)
rect(s,Inches(6.9),Inches(3.0),Inches(5.5),Inches(2.95),LIGHT)
txt(s,Inches(7.2),Inches(3.25),Inches(5),Inches(0.5),[[("RUTA DE UN MES",15,PURPLE,True)]])
txt(s,Inches(7.2),Inches(3.72),Inches(5),Inches(0.7),[[("≈ 1 mes, con profundidad",23,NAVY,True)]])
txt(s,Inches(7.2),Inches(4.5),Inches(4.9),Inches(1.3),[[("Para quienes prefieren más tiempo, afianzar conceptos y llegar al examen con mayor dominio.",13,GREY,False)]],ls=1.18)
# banda de acceso + acompañamiento
rect(s,Inches(0.9),Inches(6.12),Inches(11.5),Inches(0.62),LIGHT)
txt(s,Inches(1.1),Inches(6.18),Inches(11.1),Inches(0.5),
    [[("Acceso: ",12,PURPLE,True),
      ("1 mes en la plataforma para estudiar + 2 semanas para practicar y presentar el examen.  Extensión: $50 / 2 semanas.   ",12,NAVY,False),
      ("Acompañamiento directo del instructor de la academia.",12,PURPLE,True)]],anchor=MSO_ANCHOR.MIDDLE)
footer(s,5)

# ============ 6. PRECIO Y PROMOCIÓN ============
s=slide(); bg(s,WHITE); accent(s); logo(s); kicker(s,"Precio y promoción de lanzamiento")
title(s,"Valor de mercado, precio de apoyo a nuestras agencias")
# valor tachado vs precio
rect(s,Inches(0.9),Inches(2.05),Inches(5.6),Inches(2.6),LIGHT)
txt(s,Inches(1.2),Inches(2.3),Inches(5),Inches(0.4),[[("VALOR COMPARADO CON LA COMPETENCIA",12,GREY,True)]])
txt(s,Inches(1.2),Inches(2.75),Inches(5),Inches(0.9),[[("$350 ",40,GREY,True),("promedio",16,GREY,False)]])
txt(s,Inches(1.2),Inches(3.75),Inches(5),Inches(0.7),
    [[("Lo que cuesta prepararse con opciones equivalentes del mercado.",12.5,GREY,False)]],ls=1.15)
rect(s,Inches(6.7),Inches(2.05),Inches(5.7),Inches(2.6),PURPLE)
txt(s,Inches(7.0),Inches(2.3),Inches(5),Inches(0.4),[[("PRECIO DE LANZAMIENTO — CÓDIGO PROMOCIONAL",12,WHITE,True)]])
txt(s,Inches(7.0),Inches(2.75),Inches(5),Inches(0.9),[[("$107",46,WHITE,True)]])
txt(s,Inches(7.0),Inches(3.8),Inches(5.1),Inches(0.8),
    [[("Exclusivo para agentes de Alleanza. ",12.5,WHITE,True),
      ("Cupón aplicado en TutorLMS antes de pagar en la pasarela.",12.5,LIGHT,False)]],ls=1.15)
txt(s,Inches(0.9),Inches(5.0),Inches(11.4),Inches(1.2),
    [[("Cómo funciona:  ",13,PURPLE,True),
      ("el agente entra al curso en TutorLMS  →  aplica el código promocional  →  paga en la pasarela el precio con descuento.",13,NAVY,False)]],ls=1.2)
txt(s,Inches(0.9),Inches(6.4),Inches(11.4),Inches(0.4),
    [[("Mensaje clave: ",11,PURPLE,True),("posicionar primero el propósito (aprobar) y el valor; el precio se percibe como un beneficio, no como un gasto.",11,GREY,False,True)]])
footer(s,6)

# ============ 7. SERVICIOS OPCIONALES ============
s=slide(); bg(s,WHITE); accent(s); logo(s); kicker(s,"Servicios opcionales (a-la-carte)")
title(s,"Acompañamiento extra, separado del precio del curso")
txt(s,Inches(0.9),Inches(1.85),Inches(11.3),Inches(0.8),
    [[("No incluimos esto en el precio del curso. Quien lo necesite nos contacta directamente y lo apoyamos. Lo mostramos desglosado para total transparencia:",14,NAVY,False)]],ls=1.15)
items=[("Curso de pre-licencia (Xcel)","$30","Requisito en Florida (en Texas no aplica). Generalmente en inglés."),
       ("Cita de examen (Pearson VUE)","$50","Agendamiento del examen oficial ante Pearson VUE."),
       ("Cita de huellas","Según costo","Toma de huellas requerida para el trámite (costo del proveedor)."),
       ("Asistencia administrativa","$50","Nuestro acompañamiento para gestionar y coordinar todo el proceso.")]
x=0.9; w=2.78; g=0.16; y=2.95; h=3.0
for i,(hd,pr,bd) in enumerate(items):
    cx=Inches(x+i*(w+g))
    rect(s,cx,Inches(y),Inches(w),Inches(h),LIGHT); rect(s,cx,Inches(y),Inches(w),Inches(0.09),PURPLE)
    txt(s,cx+Inches(0.2),Inches(y+0.25),Inches(w-0.4),Inches(0.8),[[(hd,13,NAVY,True)]],ls=1.05)
    txt(s,cx+Inches(0.2),Inches(y+1.15),Inches(w-0.4),Inches(0.6),[[(pr,22,PURPLE,True)]])
    txt(s,cx+Inches(0.2),Inches(y+1.85),Inches(w-0.4),Inches(1.0),[[(bd,10.5,GREY,False)]],ls=1.12)
txt(s,Inches(0.9),Inches(6.25),Inches(11.4),Inches(0.5),
    [[("Por qué separado: ",11,PURPLE,True),("mantiene el curso principal en $107 y cubre estos costos reales sin inflar el precio de entrada.",11,GREY,False,True)]])
footer(s,7)

# ============ 8. BIBLIOTECA GRABADAS ============
s=slide(); bg(s,WHITE); accent(s); logo(s); kicker(s,"Capitalizar el conocimiento")
title(s,"Las clases grabadas como formación permanente")
txt(s,Inches(0.9),Inches(1.9),Inches(11.4),Inches(0.85),
    [[("Ya tenemos muchas clases por Zoom grabadas. Propuesta: convertir ese conocimiento en una ",15,NAVY,False),
      ("biblioteca de acceso continuo",15,PURPLE,True),(" para toda la organización.",15,NAVY,False)]],ls=1.15)
y,w,h,g=3.1,3.75,2.95,0.18
card(s,Inches(0.9),Inches(y),Inches(w),Inches(h),"Qué es",
     "Catálogo on-demand de clases ya grabadas, organizado por tema y nivel para repaso permanente.")
card(s,Inches(0.9+w+g),Inches(y),Inches(w),Inches(h),"Modelo de acceso",
     "Esquema de suscripción / acceso a evaluar: beneficio para agentes activos y/o módulo complementario.")
card(s,Inches(0.9+2*(w+g)),Inches(y),Inches(w),Inches(h),"A decidir juntos",
     "¿Gratis para retener talento, de pago para sostener la plataforma, o híbrido?")
footer(s,8)

# ============ 9. INSTRUCTORES IN-HOUSE ============
s=slide(); bg(s,NAVY); rect(s,0,0,Inches(0.18),SH,PURPLE); logo(s,dark=True)
txt(s,Inches(0.9),Inches(0.62),Inches(8),Inches(0.4),[[("CRECIMIENTO Y SOSTENIBILIDAD",13,PURPLE2,True)]])
txt(s,Inches(0.9),Inches(1.0),Inches(10),Inches(1.0),[[("Instructores que nacen dentro de Alleanza",31,WHITE,True)]])
txt(s,Inches(0.9),Inches(1.95),Inches(11.4),Inches(0.85),
    [[("A futuro, tras el curso inicial: agentes que crecieron en la casa se vuelven instructores. Comparten su experiencia real y ",14.5,LIGHT,False),
      ("generan ingresos adicionales",14.5,PURPLE2,True),(".",14.5,LIGHT,False)]],ls=1.15)
items=[("Para el agente","Una vía de ingreso extra y reconocimiento por su trayectoria dentro de la casa."),
       ("Para la academia","Banco de instructores creciente y sostenibilidad de la plataforma a largo plazo."),
       ("Para la organización","Conocimiento que se queda y se multiplica internamente; cultura que se refuerza.")]
x,w,g,y,h=0.9,3.75,0.18,3.2,2.55
for i,(hd,bd) in enumerate(items):
    card(s,Inches(x+i*(w+g)),Inches(y),Inches(w),Inches(h),hd,bd,hc=PURPLE2,dark=True)
txt(s,Inches(0.9),Inches(6.05),Inches(11.4),Inches(0.5),
    [[("Beneficio mutuo: ",13,PURPLE2,True),("oportunidades para ellos, sostenibilidad para la plataforma.",13,LIGHT,False)]])
footer(s,9,dark=True)

# ============ 10. PRÓXIMOS PASOS ============
s=slide(); bg(s,WHITE); accent(s); logo(s); kicker(s,"Próximos pasos")
title(s,"Alinearnos hoy para lanzar con solidez esta semana")
txt(s,Inches(0.9),Inches(1.95),Inches(5.6),Inches(0.5),[[("EN QUÉ NECESITAMOS ALINEARNOS",13,PURPLE,True)]])
bullets(s,Inches(0.9),Inches(2.45),Inches(5.6),Inches(3.5),[
    "Mensaje y posicionamiento ante los directores.",
    "Cómo comunicar beneficios, alcance y precio del curso.",
    "Modelo de acceso a las clases grabadas.",
    "Prioridad operativa para equipos urgentes (Miami).",
],size=13.5,gap=12)
txt(s,Inches(6.9),Inches(1.95),Inches(5.5),Inches(0.5),[[("SECUENCIA",13,PURPLE,True)]])
steps=[("1","Lanzar esta semana","Curso, rutas y promoción listos."),
       ("2","Preparar el FAQ de directores","Cómo manejar, presentar y comunicar la academia."),
       ("3","Retroalimentación post-lanzamiento","Mejoras según necesidades reales de los agentes."),
       ("4","Evolución","Instructores in-house y suscripción de grabadas.")]
yy=2.45
for num,hd,bd in steps:
    rect(s,Inches(6.9),Inches(yy),Inches(0.5),Inches(0.5),PURPLE)
    txt(s,Inches(6.9),Inches(yy),Inches(0.5),Inches(0.5),[[(num,16,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(7.55),Inches(yy-0.02),Inches(4.9),Inches(0.4),[[(hd,14,NAVY,True)]])
    txt(s,Inches(7.55),Inches(yy+0.34),Inches(4.9),Inches(0.6),[[(bd,11,GREY,False)]],ls=1.05)
    yy+=0.95
footer(s,10)

out=B+"Alleanza_Academy_Presentacion.pptx"
prs.save(out); print("saved",out,"slides:",len(prs.slides._sldIdLst))
